import os
from io import BytesIO

import pandas as pd
import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

load_dotenv()

FINANCIAL_KEYWORDS = [
    "revenue",
    "sales",
    "profit",
    "margin",
    "cost",
    "expense",
    "ebitda",
    "cash",
    "income",
    "cogs",
    "opex",
    "arr",
    "mrr",
    "ltv",
    "cac",
]

ANALYSIS_TYPES = [
    "Market Entry Analysis",
    "Competitor Analysis",
    "SWOT Analysis",
    "Pricing Strategy",
    "Go-To-Market Strategy",
    "Customer Segmentation",
    "Risk Analysis",
    "Growth Opportunities",
]

GOALS = [
    "Market entry",
    "New product launch",
    "Revenue growth",
    "Cost reduction",
    "Competitor analysis",
    "Customer acquisition",
    "Pricing strategy",
]

TIMELINES = ["0-3 months", "3-6 months", "6-12 months", "12+ months"]

PAGES = [
    "Home",
    "Financial Data Review",
    "Strategic Brief & Deliverables",
]


def init_session_state():
    defaults = {
        "df": None,
        "data_insights": None,
        "trend_df": None,
        "consultant_report": None,
        "swot_report": None,
        "competitor_report": None,
        "slide_deck": None,
        "scenario_results": None,
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


@st.cache_resource
def get_client():
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        return None
    return OpenAI(
        api_key=api_key,
        base_url="https://api.groq.com/openai/v1",
    )


def generate_consulting_insights(prompt: str) -> str:
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
    )
    return response.choices[0].message.content


def create_pptx_from_text(slide_text: str) -> BytesIO:
    prs = Presentation()
    slides = slide_text.split("# Slide")

    for slide in slides:
        if not slide.strip():
            continue
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = slide.strip().split("\n")
        title = lines[0].replace(":", "").strip()
        bullets = lines[1:]
        ppt_slide.shapes.title.text = title
        ppt_slide.placeholders[1].text = "\n".join(bullets)

    file = BytesIO()
    prs.save(file)
    file.seek(0)
    return file


def rank_kpi_columns(numeric_columns: list[str]) -> list[str]:
    scored = []
    for col in numeric_columns:
        lower = col.lower()
        score = sum(1 for kw in FINANCIAL_KEYWORDS if kw in lower)
        scored.append((score, col))
    scored.sort(key=lambda x: (-x[0], x[1]))
    return [col for _, col in scored]


def compute_trend_table(df: pd.DataFrame) -> pd.DataFrame:
    numeric_columns = df.select_dtypes(include="number").columns.tolist()
    trend_results = []

    for col in numeric_columns:
        clean_col = df[col].dropna()
        if len(clean_col) < 2:
            continue
        first_value = clean_col.iloc[0]
        last_value = clean_col.iloc[-1]
        if first_value == 0:
            continue
        percent_change = ((last_value - first_value) / first_value) * 100
        trend_results.append(
            {
                "Metric": col,
                "First Value": round(first_value, 2),
                "Last Value": round(last_value, 2),
                "Change (%)": round(percent_change, 2),
            }
        )

    return pd.DataFrame(trend_results)


def render_markdown_report(content: str, download_name: str):
    if not content:
        return

    has_sections = any(line.startswith("## ") for line in content.split("\n"))

    if not has_sections:
        st.markdown(content)
    else:
        sections = []
        current_title = "Overview"
        current_lines = []

        for line in content.split("\n"):
            if line.startswith("## "):
                if current_lines:
                    sections.append((current_title, "\n".join(current_lines)))
                current_title = line.replace("##", "").strip()
                current_lines = []
            else:
                current_lines.append(line)

        if current_lines:
            sections.append((current_title, "\n".join(current_lines)))

        for title, body in sections:
            with st.expander(title, expanded=(title.lower() == "executive summary")):
                st.markdown(body.strip() or "_No content._")

    st.download_button(
        label="Download report (.md)",
        data=content,
        file_name=download_name,
        mime="text/markdown",
        use_container_width=True,
    )


def get_brief_context() -> dict:
    return {
        "business_problem": st.session_state.get("business_problem", ""),
        "industry": st.session_state.get("industry", ""),
        "market": st.session_state.get("market", ""),
        "target_customer": st.session_state.get("target_customer", ""),
        "goal": st.session_state.get("goal", GOALS[0]),
        "timeline": st.session_state.get("timeline", TIMELINES[2]),
        "budget": st.session_state.get("budget", ""),
        "analysis_type": st.session_state.get("analysis_type", []),
        "competitors": st.session_state.get("competitors", ""),
    }


def render_brief_form(compact: bool = False):
    ctx = get_brief_context()

    business_problem = st.text_area(
        "Business question or problem",
        value=ctx["business_problem"],
        placeholder="Example: Should we expand into Southeast Asia given flat domestic margins?",
        height=100 if compact else 120,
        key="business_problem",
    )

    col1, col2 = st.columns(2)
    with col1:
        st.text_input(
            "Industry",
            value=ctx["industry"],
            placeholder="Fintech, retail, SaaS…",
            key="industry",
        )
        st.text_input(
            "Target market",
            value=ctx["market"],
            placeholder="US enterprise, Philippines…",
            key="market",
        )
    with col2:
        st.text_input(
            "Target customer",
            value=ctx["target_customer"],
            placeholder="SMB finance teams, Gen Z consumers…",
            key="target_customer",
        )
        st.selectbox("Primary goal", GOALS, index=GOALS.index(ctx["goal"]) if ctx["goal"] in GOALS else 0, key="goal")

    if not compact:
        with st.expander("Advanced options", expanded=False):
            col_a, col_b = st.columns(2)
            with col_a:
                st.selectbox(
                    "Timeline",
                    TIMELINES,
                    index=TIMELINES.index(ctx["timeline"]) if ctx["timeline"] in TIMELINES else 2,
                    key="timeline",
                )
            with col_b:
                st.text_input(
                    "Budget",
                    value=ctx["budget"],
                    placeholder="$250k, unknown, bootstrapped…",
                    key="budget",
                )
            default_analysis = ctx["analysis_type"] or ["Market Entry Analysis", "SWOT Analysis"]
            st.multiselect(
                "Analysis modules to include",
                ANALYSIS_TYPES,
                default=default_analysis,
                key="analysis_type",
            )

    return business_problem


def page_home():
    col1, col2 = st.columns(2)
    with col1:
        with st.container(border=True):
            st.markdown("**Financial Data Review**")
            if st.button("Financial Data Review", use_container_width=True):
                st.session_state["nav_page"] = "Financial Data Review"
                st.session_state["sidebar_nav_page"] = "Financial Data Review"
                st.rerun()
    with col2:
        with st.container(border=True):
            st.markdown("**Strategic Brief & Deliverables**")
            if st.button("Strategic Brief & Deliverables", use_container_width=True):
                st.session_state["nav_page"] = "Strategic Brief & Deliverables"
                st.session_state["sidebar_nav_page"] = "Strategic Brief & Deliverables"
                st.rerun()


def page_data_review():
    st.subheader("Financial Data Review")
    st.caption("Upload your data to view KPIs, trend analysis, and an AI-generated analyst memo.")

    uploaded_file = st.file_uploader(
        "Upload financial or operational data (CSV or Excel)",
        type=["csv", "xlsx"],
    )

    if uploaded_file is not None:
        if uploaded_file.name.endswith(".csv"):
            st.session_state["df"] = pd.read_csv(uploaded_file)
        else:
            st.session_state["df"] = pd.read_excel(uploaded_file)
        st.session_state["data_insights"] = None
        st.session_state["trend_df"] = None

    df = st.session_state["df"]

    if df is None:
        return

    numeric_columns = df.select_dtypes(include="number").columns.tolist()
    ranked = rank_kpi_columns(numeric_columns)

    with st.expander("Raw data preview", expanded=False):
        st.dataframe(df, use_container_width=True, height=280)
        st.caption(f"{len(df):,} rows · {len(df.columns)} columns")

    if numeric_columns:
        default_kpis = ranked[:4] if ranked else numeric_columns[:4]
        kpi_cols = st.multiselect(
            "KPI columns to highlight",
            numeric_columns,
            default=default_kpis,
            max_selections=4,
        )

        if kpi_cols:
            metric_cols = st.columns(len(kpi_cols))
            for idx, col_name in enumerate(kpi_cols):
                with metric_cols[idx]:
                    latest = df[col_name].dropna()
                    value = latest.iloc[-1] if len(latest) else 0
                    st.metric(col_name, f"{value:,.2f}")

        chart_col1, chart_col2 = st.columns([2, 1])
        with chart_col1:
            chart_metric = st.selectbox(
                "Chart metric",
                numeric_columns,
                index=numeric_columns.index(ranked[0]) if ranked else 0,
            )
        with chart_col2:
            date_candidates = ["(none)"] + df.columns.tolist()
            date_col = st.selectbox("Optional date column", date_candidates)

        if date_col != "(none)":
            chart_df = df[[date_col, chart_metric]].dropna()
            chart_df = chart_df.set_index(date_col)
            st.line_chart(chart_df, use_container_width=True)
        else:
            st.line_chart(df[chart_metric], use_container_width=True)

    st.divider()

    if st.button("Run full analysis", type="primary", use_container_width=True):
        if client is None:
            st.error("Set GROQ_API_KEY in your .env file to enable AI analysis.")
            return

        st.session_state["trend_df"] = compute_trend_table(df)

        data_preview = df.head(20).to_string()
        data_summary = df.describe().to_string()
        trend_summary = (
            st.session_state["trend_df"].to_string()
            if st.session_state["trend_df"] is not None and not st.session_state["trend_df"].empty
            else "No trend data available."
        )

        consulting_prompt = f"""
You are a junior financial analyst at a consulting firm.

Analyze this dataset and respond using EXACTLY these markdown section headers:

## Executive Summary
## Key Trends
## Risks and Red Flags
## Growth Opportunities
## Cost and Efficiency Recommendations
## Suggested Next Steps

Be specific, quantitative where possible, and practical. Flag data quality issues if relevant.

Data preview:
{data_preview}

Statistical summary:
{data_summary}

Period-over-period trends (first row vs last row per metric):
{trend_summary}
"""

        with st.spinner("Running analyst review…"):
            st.session_state["data_insights"] = generate_consulting_insights(consulting_prompt)

    if st.session_state["trend_df"] is not None and not st.session_state["trend_df"].empty:
        with st.expander("Metric trends (first vs last period)", expanded=True):
            st.dataframe(st.session_state["trend_df"], use_container_width=True, hide_index=True)

    if st.session_state["data_insights"]:
        st.markdown("### Analyst memo")
        render_markdown_report(st.session_state["data_insights"], "financial_analysis_memo.md")


def page_strategic_brief():
    st.subheader("Strategic Brief & Deliverables")
    business_problem = render_brief_form(compact=False)

    if st.button("Generate consultant brief", type="primary", use_container_width=True):
        if not business_problem.strip():
            st.warning("Enter a business question or problem first.")
        elif client is None:
            st.error("Set GROQ_API_KEY in your .env file to generate reports.")
        else:
            ctx = get_brief_context()
            consulting_prompt = f"""
You are a junior management consultant preparing a client memo.

Problem:
{ctx['business_problem']}

Industry: {ctx['industry']}
Target market: {ctx['market']}
Target customer: {ctx['target_customer']}
Business goal: {ctx['goal']}
Timeline: {ctx['timeline']}
Budget: {ctx['budget']}

Include these analysis modules: {', '.join(ctx['analysis_type']) or 'General strategic review'}

Use EXACTLY these markdown section headers:

## Executive Summary
## Market Background
## Customer Analysis
## Competitor Considerations
## Risks and Challenges
## Strategic Recommendations
## Suggested Next Steps
## Open Research Questions

Keep the response practical, structured, and business-focused.
"""
            with st.spinner("Generating consultant brief…"):
                st.session_state["consultant_report"] = generate_consulting_insights(consulting_prompt)

    if st.session_state["consultant_report"]:
        st.markdown("### Consultant report")
        render_markdown_report(st.session_state["consultant_report"], "consultant_brief.md")

    st.divider()
    st.subheader("Deliverables")
    render_deliverables_section()


def render_deliverables_section():
    ctx = get_brief_context()

    tab_swot, tab_comp, tab_fin, tab_slides = st.tabs(
        ["SWOT Analysis", "Competitor Analysis", "Scenario Planning", "Slide Deck"]
    )

    with tab_swot:
        if st.button("Generate SWOT", type="primary", use_container_width=True):
            if client is None:
                st.error("Set GROQ_API_KEY in your .env file.")
            elif not ctx["business_problem"].strip():
                st.warning("Enter a business question or problem first.")
            else:
                swot_prompt = f"""
You are a junior business consultant.

Create a SWOT for:
Problem: {ctx['business_problem']}
Industry: {ctx['industry']}
Market: {ctx['market']}
Target customer: {ctx['target_customer']}
Goal: {ctx['goal']}

Use EXACTLY these headers:

## Strengths
## Weaknesses
## Opportunities
## Threats

Be specific and practical.
"""
                with st.spinner("Generating SWOT…"):
                    st.session_state["swot_report"] = generate_consulting_insights(swot_prompt)

        if st.session_state["swot_report"]:
            render_markdown_report(st.session_state["swot_report"], "swot_analysis.md")

    with tab_comp:
        st.text_input(
            "Competitors (comma-separated)",
            value=ctx["competitors"],
            placeholder="Company A, Company B, Company C",
            key="competitors",
        )
        if st.button("Generate competitor benchmark", type="primary", use_container_width=True):
            if client is None:
                st.error("Set GROQ_API_KEY in your .env file.")
            elif not st.session_state.get("competitors", "").strip():
                st.warning("Enter at least one competitor.")
            elif not ctx["business_problem"].strip():
                st.warning("Enter a business question or problem first.")
            else:
                competitor_prompt = f"""
You are a junior business consultant.

Problem: {ctx['business_problem']}
Industry: {ctx['industry']}
Market: {ctx['market']}
Target customer: {ctx['target_customer']}
Competitors: {st.session_state['competitors']}

Use these headers:

## Competitor Overview
## Positioning Comparison
## Pricing Comparison
## Strengths and Weaknesses
## Differentiation Opportunities
## Strategic Recommendation
"""
                with st.spinner("Generating benchmark…"):
                    st.session_state["competitor_report"] = generate_consulting_insights(competitor_prompt)

        if st.session_state["competitor_report"]:
            render_markdown_report(st.session_state["competitor_report"], "competitor_benchmark.md")

    with tab_fin:
        col1, col2 = st.columns(2)
        with col1:
            base_revenue = st.number_input("Current annual revenue ($)", min_value=0.0, value=100_000.0, step=10_000.0)
            growth_rate = st.slider("Revenue growth (%)", -50, 100, 10)
        with col2:
            current_costs = st.number_input("Current annual costs ($)", min_value=0.0, value=70_000.0, step=5_000.0)
            cost_increase = st.slider("Cost increase (%)", 0, 100, 10)

        if st.button("Run scenario", type="primary", use_container_width=True):
            projected_revenue = base_revenue * (1 + growth_rate / 100)
            projected_costs = current_costs * (1 + cost_increase / 100)
            projected_profit = projected_revenue - projected_costs
            projected_margin = (projected_profit / projected_revenue) * 100 if projected_revenue else 0
            st.session_state["scenario_results"] = {
                "revenue": projected_revenue,
                "costs": projected_costs,
                "profit": projected_profit,
                "margin": projected_margin,
            }

        if st.session_state["scenario_results"]:
            r = st.session_state["scenario_results"]
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("Projected revenue", f"${r['revenue']:,.0f}")
            m2.metric("Projected costs", f"${r['costs']:,.0f}")
            m3.metric("Projected profit", f"${r['profit']:,.0f}")
            m4.metric("Profit margin", f"{r['margin']:.1f}%")

    with tab_slides:
        if st.button("Generate slide deck", type="primary", use_container_width=True):
            if client is None:
                st.error("Set GROQ_API_KEY in your .env file.")
            elif not ctx["business_problem"].strip():
                st.warning("Enter a business question or problem first.")
            else:
                slide_prompt = f"""
You are a management consultant preparing a strategy deck.

Problem: {ctx['business_problem']}
Industry: {ctx['industry']}
Market: {ctx['market']}
Target customer: {ctx['target_customer']}
Goal: {ctx['goal']}
Timeline: {ctx['timeline']}
Budget: {ctx['budget']}

Format EXACTLY like this (one slide per block):

# Slide 1: Title
- bullet
- bullet

# Slide 2: Executive Summary
- bullet

Continue through market opportunity, customer analysis, competitor landscape,
SWOT, financial considerations, risks, recommendations, and next steps.
Keep bullets concise.
"""
                with st.spinner("Building slide deck…"):
                    st.session_state["slide_deck"] = generate_consulting_insights(slide_prompt)

        if st.session_state["slide_deck"]:
            with st.expander("Preview outline", expanded=False):
                st.markdown(st.session_state["slide_deck"])
            col_a, col_b = st.columns(2)
            with col_a:
                st.download_button(
                    "Download outline (.md)",
                    data=st.session_state["slide_deck"],
                    file_name="consulting_slide_deck.md",
                    mime="text/markdown",
                    use_container_width=True,
                )
            with col_b:
                pptx_file = create_pptx_from_text(st.session_state["slide_deck"])
                st.download_button(
                    "Download PowerPoint (.pptx)",
                    data=pptx_file,
                    file_name="consulting_slide_deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )


# --- App shell ---
st.set_page_config(
    page_title="Junior Analyst Copilot",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
    <style>
    div[data-testid="stMetric"] {
        border: 1px solid rgba(128, 128, 128, 0.35);
        border-radius: 8px;
        padding: 12px 16px;
    }
    div[data-testid="stMetric"] [data-testid="stMetricLabel"] p,
    div[data-testid="stMetric"] [data-testid="stMetricValue"] {
        color: inherit;
    }
    .block-container { padding-top: 2rem; }
    </style>
    """,
    unsafe_allow_html=True,
)

init_session_state()
client = get_client()

if "nav_page" not in st.session_state:
    st.session_state["nav_page"] = "Home"
if "sidebar_nav_page" not in st.session_state:
    st.session_state["sidebar_nav_page"] = st.session_state["nav_page"]

_legacy_pages = {
    "Strategic Brief": "Strategic Brief & Deliverables",
    "Deliverables": "Strategic Brief & Deliverables",
}
if st.session_state["nav_page"] in _legacy_pages:
    st.session_state["nav_page"] = _legacy_pages[st.session_state["nav_page"]]
if st.session_state["sidebar_nav_page"] in _legacy_pages:
    st.session_state["sidebar_nav_page"] = _legacy_pages[st.session_state["sidebar_nav_page"]]


def sync_from_sidebar():
    st.session_state["nav_page"] = st.session_state["sidebar_nav_page"]


st.session_state["sidebar_nav_page"] = st.session_state["nav_page"]

with st.sidebar:
    st.title("Junior Analyst Copilot")

    st.radio(
        "Navigate",
        PAGES,
        key="sidebar_nav_page",
        on_change=sync_from_sidebar,
        label_visibility="collapsed",
    )

    st.divider()
    if client is None:
        st.error("GROQ_API_KEY not found in .env")
    else:
        st.success("Groq API connected")

    if st.button("Clear session", use_container_width=True):
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()

st.title("Junior Analyst Copilot")

page = st.session_state["nav_page"]

if page == "Home":
    page_home()
elif page == "Financial Data Review":
    page_data_review()
elif page == "Strategic Brief & Deliverables":
    page_strategic_brief()
